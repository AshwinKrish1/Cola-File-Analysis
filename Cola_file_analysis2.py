import streamlit as st
import os
import csv
from PIL import Image
import pytesseract
import fitz  # PyMuPDF for PDF
from pptx import Presentation
import pandas as pd

# Set path to Tesseract executable if not in system PATH
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Define keywords to search for (case-insensitive)
keywords = ['cola', 'cost of living allowance']

# Supported file extensions
image_extensions = ['.png', '.jpg', '.jpeg']
text_extensions = ['.txt', '.log', '.csv']
pdf_extensions = ['.pdf']
ppt_extensions = ['.ppt', '.pptx']

# Function to check for keywords in text (case-insensitive)
def contains_keywords(text):
    text_lower = text.lower()
    return any(keyword in text_lower for keyword in keywords)

# Function to process files in a folder
def process_folder(parent_folder):
    results = []
    for root, dirs, files in os.walk(parent_folder):
        for filename in files:
            file_path = os.path.join(root, filename)
            ext = os.path.splitext(filename)[1].lower()
            found = False
            try:
                if ext in text_extensions:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                        found = contains_keywords(content)
                elif ext in image_extensions:
                    text = pytesseract.image_to_string(Image.open(file_path))
                    found = contains_keywords(text)
                elif ext in pdf_extensions:
                    doc = fitz.open(file_path)
                    pdf_text = ""
                    for page in doc:
                        pdf_text += page.get_text()
                    found = contains_keywords(pdf_text)
                elif ext in ppt_extensions:
                    prs = Presentation(file_path)
                    ppt_text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                ppt_text += shape.text + "\n"
                    found = contains_keywords(ppt_text)
                results.append((root, filename, 'Yes' if found else 'No'))
            except Exception as e:
                results.append((root, filename, f'Error: {str(e)}'))
    return results

# Streamlit UI
st.set_page_config(page_title="COLA Analysis", layout="wide")
st.title("COLA Analysis")
folder_path = st.text_input("Enter the parent folder path:")

if folder_path:
    if os.path.isdir(folder_path):
        st.info("Processing files. Please wait...")
        data = process_folder(folder_path)
        df = pd.DataFrame(data, columns=['Folder Path', 'Filename', 'Contains Keywords'])
        st.dataframe(df, use_container_width=True)

        # Download CSV
        csv_file = "keyword_search_results.csv"
        df.to_csv(csv_file, index=False)
        with open(csv_file, "rb") as f:
            st.download_button("Download Results as CSV", f, file_name=csv_file)
    else:
        st.error("Invalid folder path. Please enter a valid directory.")